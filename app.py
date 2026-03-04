"""
Walmart Image → PPT 자동화 (Streamlit)
1. 이미지 폴더(zip) 또는 파일 업로드
2. PPT 템플릿 업로드
3. 자동 생성 → 다운로드
"""

import streamlit as st
import tempfile
import zipfile
import os
import re
import io
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from PIL import Image


# ── 파일명 파싱 ──────────────────────────────────────────────

def parse_image_filename(filename):
    name = os.path.splitext(os.path.basename(filename))[0]

    if name.endswith("_Front"):
        view, name = "Front", name[:-6]
    elif name.endswith("_Back"):
        view, name = "Back", name[:-5]
    else:
        view = "Side"

    m = re.search(r"(S\d+\s+\d+[_ ]D\d+)", name)
    design_key = name[m.start():] if m else name

    return design_key, _extract_colorway(design_key), view


def _extract_colorway(key):
    for pattern in [
        r"_CW\d+_(.+?)_WM$",
        r"_([^_]+)_WM$",
        r"\((.+?)\)",
    ]:
        m = re.search(pattern, key)
        if m:
            return m.group(1)
    return ""


def group_images(file_map):
    """
    file_map: {filename: filepath} → 디자인별 그룹핑
    Returns: {design_key: {artwork, colorway, front, back}}
    """
    groups = {}
    for filename in sorted(file_map.keys()):
        ext = os.path.splitext(filename)[1].lower()
        if ext not in (".png", ".jpg", ".jpeg"):
            continue

        design_key, colorway, view = parse_image_filename(filename)

        if design_key not in groups:
            groups[design_key] = {"colorway": colorway, "artwork": design_key}

        if view in ("Front", "Back"):
            groups[design_key][view.lower()] = file_map[filename]

    return {k: v for k, v in groups.items() if "front" in v and "back" in v}


# ── 슬라이드 복제 ────────────────────────────────────────────

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def duplicate_slide(prs, source_slide):
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    spTree = new_slide.shapes._spTree
    for child in list(spTree):
        if etree.QName(child.tag).localname not in ("nvGrpSpPr", "grpSpPr"):
            spTree.remove(child)

    rid_map = {}
    for rel in source_slide.part.rels.values():
        if "image" in rel.reltype:
            new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rel.rId] = new_rid

    for child in source_slide.shapes._spTree:
        if etree.QName(child.tag).localname in ("nvGrpSpPr", "grpSpPr"):
            continue
        new_child = deepcopy(child)
        for elem in new_child.iter():
            for attr in (f"{{{R_NS}}}embed", f"{{{R_NS}}}link"):
                old_val = elem.get(attr)
                if old_val and old_val in rid_map:
                    elem.set(attr, rid_map[old_val])
        spTree.append(new_child)

    return new_slide


# ── 슬라이드 내용 채우기 ─────────────────────────────────────

def _calc_image_rect(img_path, slide_width, img_top, img_bottom, img_max_width):
    with Image.open(img_path) as img:
        px_w, px_h = img.size

    ratio = px_h / px_w
    max_h = img_bottom - img_top

    if img_max_width * ratio <= max_h:
        w = img_max_width
        h = int(w * ratio)
    else:
        h = max_h
        w = int(h / ratio)

    left = (slide_width - w) // 2
    top = img_top + (max_h - h) // 2
    return left, top, w, h


def fill_slide(slide, image_path, artwork, colorway, slide_width):
    info_top = None
    title_bottom = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if "Artwork" in text or "Colorway" in text:
            info_top = shape.top
        else:
            title_bottom = max(title_bottom, shape.top + shape.height)

    if info_top is None:
        info_top = int(slide_width * 1.2)

    img_area_top = title_bottom + 200000
    img_area_bottom = info_top - 100000
    img_max_width = int(slide_width * 0.6)

    left, top, w, h = _calc_image_rect(
        image_path, slide_width, img_area_top, img_area_bottom, img_max_width
    )
    slide.shapes.add_picture(image_path, left, top, w, h)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            text = p.text.strip()
            if text == "Artwork:" or text.startswith("Artwork:"):
                for run in p.runs:
                    if "Artwork" in run.text:
                        run.text = f"Artwork: {artwork}"
                        break
            if text.startswith("Colorway"):
                for run in p.runs:
                    if "Colorway" in run.text:
                        run.text = f"Colorway : {colorway}" if colorway else "Colorway : N/A"
                        break


# ── PPT 생성 ─────────────────────────────────────────────────

def generate_ppt(template_bytes, file_map, progress_bar=None):
    groups = group_images(file_map)
    if not groups:
        raise ValueError("Front/Back 이미지 쌍을 찾을 수 없습니다.\n파일명에 _Front / _Back 접미사가 필요합니다.")

    prs = Presentation(io.BytesIO(template_bytes))
    template_slide = prs.slides[0]
    slide_width = prs.slide_width

    designs = list(groups.items())
    total = len(designs) * 2

    all_slides = [template_slide]
    for _ in range(total - 1):
        all_slides.append(duplicate_slide(prs, template_slide))

    for i, (key, g) in enumerate(designs):
        fill_slide(all_slides[i * 2], g["front"], g["artwork"], g["colorway"], slide_width)
        if progress_bar:
            progress_bar.progress((i * 2 + 1) / total)

        fill_slide(all_slides[i * 2 + 1], g["back"], g["artwork"], g["colorway"], slide_width)
        if progress_bar:
            progress_bar.progress((i * 2 + 2) / total)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output, len(designs)


# ── Streamlit UI ─────────────────────────────────────────────

st.set_page_config(page_title="Walmart Image → PPT", page_icon="📎", layout="centered")
st.title("Walmart Image → PPT 자동화")

# 1. 이미지 업로드
st.subheader("1. 이미지 업로드")
upload_mode = st.radio("업로드 방식", ["이미지 파일 직접 선택", "ZIP 파일 업로드"], horizontal=True)

image_files = None
if upload_mode == "이미지 파일 직접 선택":
    image_files = st.file_uploader(
        "Front / Back 이미지를 모두 선택하세요",
        type=["png", "jpg", "jpeg"],
        accept_multiple_files=True,
    )
else:
    zip_file = st.file_uploader("이미지가 들어있는 ZIP 파일", type=["zip"])
    if zip_file:
        image_files = zip_file  # 나중에 처리

# 2. PPT 템플릿
st.subheader("2. PPT 템플릿 업로드")
template_file = st.file_uploader("빈 PPT 템플릿 (.pptx)", type=["pptx"])

# 3. 생성
st.divider()

if st.button("PPT 생성", type="primary", use_container_width=True):
    if not image_files:
        st.error("이미지를 업로드하세요.")
    elif not template_file:
        st.error("PPT 템플릿을 업로드하세요.")
    else:
        with st.spinner("PPT 생성 중..."):
            tmp_dir = tempfile.mkdtemp()
            file_map = {}

            try:
                # 이미지 파일을 임시 디렉토리에 저장
                if upload_mode == "ZIP 파일 업로드":
                    with zipfile.ZipFile(io.BytesIO(image_files.read())) as zf:
                        for name in zf.namelist():
                            if name.startswith("__MACOSX") or name.startswith("."):
                                continue
                            basename = os.path.basename(name)
                            ext = os.path.splitext(basename)[1].lower()
                            if ext in (".png", ".jpg", ".jpeg"):
                                tmp_path = os.path.join(tmp_dir, basename)
                                with open(tmp_path, "wb") as f:
                                    f.write(zf.read(name))
                                file_map[basename] = tmp_path
                else:
                    for uf in image_files:
                        tmp_path = os.path.join(tmp_dir, uf.name)
                        with open(tmp_path, "wb") as f:
                            f.write(uf.read())
                        file_map[uf.name] = tmp_path

                if not file_map:
                    st.error("업로드된 이미지가 없습니다.")
                else:
                    progress = st.progress(0)
                    template_bytes = template_file.read()

                    output, count = generate_ppt(template_bytes, file_map, progress)
                    progress.progress(1.0)

                    st.success(f"{count}개 디자인 × 2 = {count * 2}개 슬라이드 생성 완료!")
                    st.download_button(
                        label="📥 PPT 다운로드",
                        data=output,
                        file_name=f"{os.path.splitext(template_file.name)[0]}_완료.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                    )

            except Exception as e:
                st.error(f"오류: {e}")
            finally:
                # 임시 파일 정리
                import shutil
                shutil.rmtree(tmp_dir, ignore_errors=True)
